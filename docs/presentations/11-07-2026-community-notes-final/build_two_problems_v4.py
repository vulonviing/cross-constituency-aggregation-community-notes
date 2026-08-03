from __future__ import annotations

import hashlib
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


HERE = Path(__file__).resolve().parent
SOURCE = HERE / "community-notes-final-presentation-v2.pptx"
OUTPUT = HERE / "community-notes-final-presentation-v4.pptx"
SOURCE_SHA256 = "7f1aa1bb21f8aad811e03076cd624c5fee79a5a949617fed98149c5b1030b992"
TARGETS = {
    "ppt/slides/slide7.xml": "ppt/slides/slide1.xml",
    "ppt/slides/slide8.xml": "ppt/slides/slide2.xml",
}

W, H = 13.333, 7.5
FONT = "Helvetica Neue"
WHITE = "FFFFFF"
INK = "111318"
MUTED = "73777F"
LIGHT = "DADDE2"
FAINT = "ECEEF1"
BLUE = "3977F6"
CORAL = "FF705B"
GREEN = "1FA873"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def textbox(slide, text, x, y, w, h, size=24, color=INK, bold=False,
            align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear(); frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align; paragraph.space_after = Pt(0)
    run = paragraph.add_run(); run.text = text
    run.font.name = FONT; run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def bullet(slide, text, x, y, w, size=15, color=INK, dot=GREEN, bold=False,
           align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.40))
    frame = box.text_frame
    frame.clear(); frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]; paragraph.alignment = align
    r1 = paragraph.add_run(); r1.text = "• "; r1.font.name = FONT
    r1.font.size = Pt(size); r1.font.bold = True; r1.font.color.rgb = rgb(dot)
    r2 = paragraph.add_run(); r2.text = text; r2.font.name = FONT
    r2.font.size = Pt(size); r2.font.bold = bold; r2.font.color.rgb = rgb(color)
    return box


def line(slide, x1, y1, x2, y2, color=INK, width=1.5):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    connector.line.color.rgb = rgb(color); connector.line.width = Pt(width)
    return connector


def circle(slide, x, y, diameter, fill=INK, outline=None, width=0.8):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(diameter), Inches(diameter)
    )
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(outline or fill); shape.line.width = Pt(width)
    return shape


def outline_round_rect(slide, x, y, w, h, color=INK, width=1.1):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.background(); shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)
    return shape


def base_slide(presentation):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(WHITE)
    return slide


def title(slide, kicker, text, number):
    textbox(slide, kicker, 0.66, 0.46, 4.8, 0.25, 10, MUTED, True)
    textbox(slide, text, 0.66, 0.86, 11.5, 0.62, 30, INK, True)
    textbox(slide, number, 12.12, 0.50, 0.48, 0.20, 9, MUTED, True, PP_ALIGN.RIGHT)


def source(slide, text):
    textbox(slide, f"• {text}", 0.68, 7.04, 11.9, 0.18, 8.2, MUTED)


def user(slide, x, y, scale=1.0, color=INK, raised=False):
    head = 0.20 * scale
    circle(slide, x, y, head, color)
    cx = x + head / 2
    line(slide, cx, y + head, cx, y + 0.53 * scale, color, 1.5)
    line(slide, x - 0.04 * scale, y + 0.34 * scale,
         x + 0.24 * scale, y + 0.34 * scale, color, 1.5)
    if raised:
        line(slide, cx + 0.08 * scale, y + 0.34 * scale,
             cx + 0.23 * scale, y + 0.05 * scale, color, 1.8)
    line(slide, cx, y + 0.53 * scale, x - 0.03 * scale, y + 0.74 * scale, color, 1.5)
    line(slide, cx, y + 0.53 * scale, x + 0.23 * scale, y + 0.74 * scale, color, 1.5)


def note(slide, x, y, w=1.35, h=1.72, color=INK):
    outline_round_rect(slide, x, y, w, h, color, 1.1)
    line(slide, x + 0.24, y + 0.48, x + w - 0.24, y + 0.48, color, 1.2)
    line(slide, x + 0.24, y + 0.82, x + w - 0.38, y + 0.82, color, 1.2)
    line(slide, x + 0.24, y + 1.16, x + w - 0.28, y + 1.16, color, 1.2)


def helpful_mark(slide, x, y, color=GREEN):
    circle(slide, x, y, 0.22, color)
    line(slide, x + 0.045, y + 0.12, x + 0.095, y + 0.17, WHITE, 1.3)
    line(slide, x + 0.095, y + 0.17, x + 0.185, y + 0.06, WHITE, 1.3)


def build_problem1(slide):
    title(slide, "CORE PROBLEM 1", "It reads the rater", "04")
    line(slide, 6.67, 1.78, 6.67, 5.62, LIGHT, 1.0)

    # Amigo: expected partisan support is absorbed by viewpoint compatibility.
    textbox(slide, "AMIGO", 0.90, 1.72, 1.8, 0.24, 10, CORAL, True)
    user(slide, 1.35, 2.43, 1.25, CORAL)
    note(slide, 3.62, 2.20, 1.48, 1.88, CORAL)
    helpful_mark(slide, 2.55, 2.83)
    textbox(slide, "HELPFUL", 2.30, 3.18, 0.78, 0.20, 9, GREEN, True, PP_ALIGN.CENTER)
    line(slide, 2.01, 2.94, 2.46, 2.94, CORAL, 1.4)
    line(slide, 2.82, 2.94, 3.58, 2.94, CORAL, 1.4)
    textbox(slide, "EXPECTED", 1.15, 4.47, 1.50, 0.25, 11, MUTED, True)
    textbox(slide, "fᵤ · fₙ", 2.72, 4.30, 1.30, 0.42, 23, CORAL, True, PP_ALIGN.CENTER)
    textbox(slide, "iₙ  —", 4.35, 4.32, 1.22, 0.38, 21, MUTED, True, PP_ALIGN.CENTER)
    line(slide, 2.45, 4.55, 2.70, 4.55, LIGHT, 1.3)
    line(slide, 4.00, 4.55, 4.30, 4.55, LIGHT, 1.3)
    bullet(slide, "Predictable partisan support", 1.08, 5.22, 3.5, 14, MUTED, CORAL)

    # Enemies: the same Helpful action receives extra note-level credit.
    textbox(slide, "ENEMIES SHAKE HANDS", 7.02, 1.72, 2.7, 0.24, 10, GREEN, True)
    user(slide, 7.55, 2.43, 1.18, BLUE)
    note(slide, 9.58, 2.20, 1.48, 1.88, INK)
    user(slide, 11.75, 2.43, 1.18, CORAL)
    helpful_mark(slide, 8.72, 2.83)
    helpful_mark(slide, 11.13, 2.83)
    line(slide, 8.15, 2.94, 8.62, 2.94, BLUE, 1.4)
    line(slide, 8.98, 2.94, 9.54, 2.94, BLUE, 1.4)
    line(slide, 11.17, 2.94, 11.66, 2.94, CORAL, 1.4)
    textbox(slide, "UNEXPECTED", 7.38, 4.47, 1.78, 0.25, 11, GREEN, True)
    textbox(slide, "fᵤ · fₙ ≈ 0", 9.14, 4.30, 1.85, 0.42, 21, MUTED, True, PP_ALIGN.CENTER)
    textbox(slide, "iₙ  ↑", 11.20, 4.27, 1.00, 0.42, 24, GREEN, True, PP_ALIGN.CENTER)
    line(slide, 8.88, 4.55, 9.12, 4.55, LIGHT, 1.3)
    line(slide, 11.00, 4.55, 11.18, 4.55, LIGHT, 1.3)
    bullet(slide, "Cross-viewpoint surprise", 7.42, 5.22, 3.4, 14, MUTED, GREEN)

    line(slide, 0.82, 5.92, 12.50, 5.92, LIGHT, 1.1)
    bullet(slide, "The same click receives different credit", 0.92, 6.20, 4.5,
           16, INK, BLUE, True)
    bullet(slide, "The model asks who voted", 7.42, 6.20, 4.1,
           16, INK, CORAL, True)
    source(slide, "Wojcik et al. (2022); Buterin (2023); X Community Notes (2026)")


def build_problem2(slide):
    # Split the kicker label and number to avoid a QuickLook/Keynote clipping
    # quirk observed when the full phrase is stored in one compact text box.
    textbox(slide, "CORE PROBLEM", 0.66, 0.46, 1.30, 0.25, 10, MUTED, True)
    textbox(slide, "2", 1.94, 0.46, 0.20, 0.25, 10, MUTED, True)
    textbox(slide, "The same hands draw the map", 0.66, 0.86, 11.5, 0.62, 30, INK, True)
    textbox(slide, "05", 12.12, 0.50, 0.48, 0.20, 9, MUTED, True, PP_ALIGN.RIGHT)
    textbox(slide, "CLASSROOM METAPHOR", 0.88, 1.70, 2.3, 0.22, 10, MUTED, True)

    active_positions = {(1, 1), (4, 0), (5, 2)}
    active_centers = []
    for row in range(4):
        for col in range(7):
            x = 0.95 + col * 0.68 + (0.18 if row % 2 else 0)
            y = 2.15 + row * 0.80
            active = (col, row) in active_positions
            color = CORAL if active else MUTED
            user(slide, x, y, 0.72, color, raised=active)
            if active:
                active_centers.append((x + 0.15, y + 0.04))
                for tick in range(4):
                    line(slide, x + 0.43 + tick * 0.07, y + 0.03,
                         x + 0.43 + tick * 0.07, y + 0.14, CORAL, 1.2)

    # Repeated answers flow into the learned viewpoint map.
    map_x1, map_x2, map_y = 7.10, 12.22, 3.55
    line(slide, map_x1, map_y, map_x2, map_y, LIGHT, 2.0)
    for px, py in active_centers:
        line(slide, px + 0.30, py + 0.20, map_x1, map_y, CORAL, 1.0)
        line(slide, px + 0.34, py + 0.27, map_x1, map_y + 0.06, CORAL, 0.8)
    # Small ordinary-user signal remains visible but faint.
    line(slide, 4.76, 4.70, map_x1, map_y + 0.14, LIGHT, 0.7)
    line(slide, 1.98, 3.10, map_x1, map_y - 0.12, LIGHT, 0.7)

    points = [
        (7.42, 3.35, BLUE, 0.10), (7.88, 3.78, BLUE, 0.09),
        (8.34, 3.30, BLUE, 0.11), (9.02, 3.68, MUTED, 0.08),
        (9.72, 3.20, CORAL, 0.14), (10.16, 3.73, CORAL, 0.13),
        (10.68, 3.16, CORAL, 0.16), (11.20, 3.62, CORAL, 0.14),
        (11.74, 3.26, CORAL, 0.15),
    ]
    for x, y, color, diameter in points:
        circle(slide, x, y, diameter, color)
    textbox(slide, "VIEWPOINT MAP", 7.08, 4.18, 2.0, 0.22, 10, MUTED, True)
    textbox(slide, "64", 8.12, 4.72, 1.00, 0.58, 34, CORAL, True)
    textbox(slide, "• power raters", 9.10, 4.91, 1.65, 0.25, 12, MUTED, True)
    textbox(slide, "11×", 10.72, 4.72, 1.12, 0.58, 34, INK, True)
    textbox(slide, "• more active", 11.63, 4.91, 1.15, 0.25, 12, MUTED, True)

    line(slide, 0.82, 5.75, 12.50, 5.75, LIGHT, 1.1)
    bullet(slide, "A few users supply most of the map", 0.92, 6.10, 4.8,
           16, INK, CORAL, True)
    bullet(slide, "The active core becomes ‘the community’", 7.38, 6.10, 4.8,
           16, INK, BLUE, True)
    source(slide, "Nudo et al. (2026); authors’ 200k analysis")


def build_replacements(path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(W); presentation.slide_height = Inches(H)
    slide1 = base_slide(presentation); build_problem1(slide1)
    slide2 = base_slide(presentation); build_problem2(slide2)
    presentation.save(path)


def build_single_preview(path: Path, number: int) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(W); presentation.slide_height = Inches(H)
    slide = base_slide(presentation)
    (build_problem1 if number == 1 else build_problem2)(slide)
    presentation.save(path)


def replace_slide_xml(replacement_pptx: Path) -> None:
    actual_hash = sha256(SOURCE)
    if actual_hash != SOURCE_SHA256:
        raise SystemExit(
            "Source V2 changed; refusing to overwrite manual edits. "
            f"Expected {SOURCE_SHA256}, got {actual_hash}."
        )
    with zipfile.ZipFile(replacement_pptx, "r") as replacement_zip:
        replacement_xml = {
            target: replacement_zip.read(source)
            for target, source in TARGETS.items()
        }
    temporary = OUTPUT.with_suffix(".tmp.pptx")
    if temporary.exists():
        temporary.unlink()
    with zipfile.ZipFile(SOURCE, "r") as source_zip, zipfile.ZipFile(temporary, "w") as output_zip:
        for info in source_zip.infolist():
            output_zip.writestr(info, replacement_xml.get(info.filename, source_zip.read(info.filename)))
    temporary.replace(OUTPUT)


def main() -> None:
    with tempfile.TemporaryDirectory() as tmpdir:
        replacement = Path(tmpdir) / "two-problems.pptx"
        build_replacements(replacement)
        replace_slide_xml(replacement)
    print(f"Source V2 preserved: {sha256(SOURCE)}")
    print(f"Saved V4: {OUTPUT}")


if __name__ == "__main__":
    main()
