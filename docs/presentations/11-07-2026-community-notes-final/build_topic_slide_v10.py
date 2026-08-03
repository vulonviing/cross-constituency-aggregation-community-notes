from __future__ import annotations

import hashlib
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from build_two_problems_v4 import (
    BLUE,
    CORAL,
    GREEN,
    INK,
    LIGHT,
    MUTED,
    W,
    H,
    base_slide,
    bullet,
    line,
    source,
    textbox,
    title,
)


HERE = Path(__file__).resolve().parent
ROOT = HERE.parents[2]
SOURCE = HERE / "community-notes-final-presentation-v9.pptx"
OUTPUT = HERE / "community-notes-final-presentation-v10.pptx"
FIGURE = ROOT / "figures/script_figures/cn-topic-signatures.png"
SOURCE_SHA256 = "4699a7a84a9a10d0d1e7429ddaee3b72faf2c8d4ba678d052149273930d34991"
FIGURE_SHA256 = "a90dd89fc6c240f8a3952d5efe288e1d0569ed5b33affb75a817f18d89a04494"


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def topic_slide(slide) -> None:
    title(slide, "TOPIC MODELING", "Disagreement changes by topic", "10")

    # Preserve the paper figure exactly: proportional scaling, no crop or edit.
    slide.shapes.add_picture(str(FIGURE), Inches(0.68), Inches(1.56), width=Inches(7.18))

    textbox(slide, "HOW TO READ", 8.28, 1.63, 2.10, 0.23, 10, MUTED, True)
    bullet(slide, "x-position = approval rate", 8.28, 2.05, 3.75, 13, MUTED, BLUE)
    bullet(slide, "bubble size = note count", 8.28, 2.50, 3.75, 13, MUTED, CORAL)
    bullet(slide, "color = recovered constituency", 8.28, 2.95, 3.92, 13, MUTED, GREEN)

    line(slide, 8.28, 3.54, 12.28, 3.54, LIGHT, 1.0)
    textbox(slide, "KEY RESULT", 8.28, 3.87, 2.10, 0.23, 10, GREEN, True)
    bullet(slide, "Approval leadership flips across topics", 8.28, 4.29, 3.92, 13.2, INK, GREEN, True)
    bullet(slide, "Constituencies are content-dependent", 8.28, 4.84, 3.92, 13.2, INK, GREEN, True)

    line(slide, 8.28, 5.52, 12.28, 5.52, LIGHT, 1.0)
    bullet(slide, "The clusters are not simply ‘strict’ versus ‘lenient’",
           8.28, 5.82, 3.98, 13.2, INK, CORAL, True)
    source(slide, "Authors’ BERTopic analysis · Method-B clusters")


def build_generated(path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(W)
    presentation.slide_height = Inches(H)
    topic_slide(base_slide(presentation))
    presentation.save(path)


def replace_once(payload: bytes, old: bytes, new: bytes, label: str) -> bytes:
    count = payload.count(old)
    if count != 1:
        raise RuntimeError(f"Expected one {label} marker, found {count}")
    return payload.replace(old, new, 1)


def renumber_slide(payload: bytes, old: str, new: str, slide_no: int) -> bytes:
    old_marker = f'<a:t>{old}</a:t>'.encode()
    new_marker = f'<a:t>{new}</a:t>'.encode()
    count = payload.count(old_marker)
    if count != 1:
        raise RuntimeError(
            f"Expected one visible number {old!r} on package slide {slide_no}, found {count}"
        )
    return payload.replace(old_marker, new_marker, 1)


def assemble(generated_path: Path) -> None:
    if sha256(SOURCE) != SOURCE_SHA256:
        raise SystemExit("V9 changed; refusing to overwrite the manually edited source.")
    if sha256(FIGURE) != FIGURE_SHA256:
        raise SystemExit("Topic figure changed; refusing to embed an unreviewed image.")

    with zipfile.ZipFile(generated_path, "r") as generated:
        new_slide = generated.read("ppt/slides/slide1.xml")
        new_rels = generated.read("ppt/slides/_rels/slide1.xml.rels")
        new_image = generated.read("ppt/media/image1.png")
        new_rels = new_rels.replace(b"../media/image1.png", b"../media/image3.png")

    with zipfile.ZipFile(SOURCE, "r") as source_zip:
        replacements: dict[str, bytes] = {}

        # Only the displayed section number changes on the following slides.
        for slide_no in range(24, 32):
            old_number = "10" if slide_no == 24 else ("11" if slide_no <= 28 else "12")
            new_number = "11" if slide_no == 24 else ("12" if slide_no <= 28 else "13")
            name = f"ppt/slides/slide{slide_no}.xml"
            replacements[name] = renumber_slide(
                source_zip.read(name), old_number, new_number, slide_no
            )

        presentation_xml = source_zip.read("ppt/presentation.xml")
        presentation_xml = replace_once(
            presentation_xml,
            b'<p:sldId id="279" r:id="rId31"/>',
            b'<p:sldId id="287" r:id="rId39"/><p:sldId id="279" r:id="rId31"/>',
            "slide 24 insertion point",
        )
        replacements["ppt/presentation.xml"] = presentation_xml

        presentation_rels = source_zip.read("ppt/_rels/presentation.xml.rels")
        relationship = (
            b'<Relationship Id="rId39" '
            b'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" '
            b'Target="slides/slide32.xml"/>'
        )
        replacements["ppt/_rels/presentation.xml.rels"] = replace_once(
            presentation_rels,
            b"</Relationships>",
            relationship + b"</Relationships>",
            "presentation relationship",
        )

        content_types = source_zip.read("[Content_Types].xml")
        override = (
            b'<Override PartName="/ppt/slides/slide32.xml" '
            b'ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'
        )
        replacements["[Content_Types].xml"] = replace_once(
            content_types, b"</Types>", override + b"</Types>", "content type"
        )

        additions = {
            "ppt/slides/slide32.xml": new_slide,
            "ppt/slides/_rels/slide32.xml.rels": new_rels,
            "ppt/media/image3.png": new_image,
        }

        temporary = OUTPUT.with_suffix(".tmp.pptx")
        if temporary.exists():
            temporary.unlink()
        with zipfile.ZipFile(temporary, "w") as output_zip:
            for info in source_zip.infolist():
                output_zip.writestr(
                    info, replacements.get(info.filename, source_zip.read(info.filename))
                )
            for name, payload in additions.items():
                output_zip.writestr(name, payload)
        temporary.replace(OUTPUT)


def main() -> None:
    with tempfile.TemporaryDirectory() as directory:
        generated = Path(directory) / "topic.pptx"
        build_generated(generated)
        assemble(generated)
    print(f"Created {OUTPUT}")


if __name__ == "__main__":
    main()
