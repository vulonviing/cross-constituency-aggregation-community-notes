from __future__ import annotations

import hashlib
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Inches, Pt


HERE = Path(__file__).resolve().parent
SOURCE = HERE / "community-notes-final-presentation-v10.pptx"
OUTPUT = HERE / "community-notes-final-presentation-v11.pptx"
SOURCE_SHA256 = "c2450c224f84eaa363548ff3464693118e335bae434bb00cf7769411e3d6c0ca"

LIGHT = RGBColor.from_string("DADDE2")
CORAL = RGBColor.from_string("FF705B")

# Seventeen new lines join the seven already present on Slide 6.
# Coordinates are one-based conceptually but stored as zero-based row/column.
NEW_VOTERS = [
    (0, 2), (0, 7), (0, 11),
    (1, 4), (1, 10),
    (2, 1), (2, 13),
    (3, 3), (3, 8),
    (4, 0), (4, 10),
    (5, 5), (5, 14),
    (6, 2), (6, 8),
    (7, 12),
    (8, 4),
]
DOMINANT_VOTERS = {(0, 11), (1, 4), (3, 3), (4, 10), (5, 5), (6, 8)}


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def grid_dot(slide, row: int, col: int):
    target = 10 + row * 15 + col
    name = f"Oval {target}"
    matches = [shape for shape in slide.shapes if shape.name == name]
    if len(matches) != 1:
        raise RuntimeError(f"Expected one grid dot {name}, found {len(matches)}")
    return matches[0]


def add_vote_lines(presentation: Presentation) -> None:
    slide = presentation.slides[5]
    note_x = Inches(10.42)
    note_y = Inches(3.38)
    for row, col in NEW_VOTERS:
        dot = grid_dot(slide, row, col)
        start_x = dot.left + dot.width // 2
        start_y = dot.top + dot.height // 2
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, start_x, start_y, note_x, note_y
        )
        connector.name = f"VoteLine r{row + 1}c{col + 1}"
        connector.line.color.rgb = LIGHT
        connector.line.width = Pt(0.70)


def add_dominance_emphasis(presentation: Presentation) -> None:
    slide = presentation.slides[5]
    for row, col in DOMINANT_VOTERS:
        dot = grid_dot(slide, row, col)
        dot.fill.solid()
        dot.fill.fore_color.rgb = CORAL
        dot.line.color.rgb = CORAL

        line_name = f"VoteLine r{row + 1}c{col + 1}"
        matches = [shape for shape in slide.shapes if shape.name == line_name]
        if len(matches) != 1:
            raise RuntimeError(f"Expected one dominant line {line_name}, found {len(matches)}")
        matches[0].line.color.rgb = CORAL
        matches[0].line.width = Pt(2.20)


def build_variants(base_path: Path, emphasis_path: Path) -> None:
    base = Presentation(SOURCE)
    add_vote_lines(base)
    base.save(base_path)

    emphasis = Presentation(base_path)
    add_dominance_emphasis(emphasis)
    emphasis.save(emphasis_path)


def replace_once(payload: bytes, old: bytes, new: bytes, label: str) -> bytes:
    count = payload.count(old)
    if count != 1:
        raise RuntimeError(f"Expected one {label} marker, found {count}")
    return payload.replace(old, new, 1)


def assemble(base_path: Path, emphasis_path: Path) -> None:
    actual_hash = sha256(SOURCE)
    if actual_hash != SOURCE_SHA256:
        raise SystemExit(
            "V10 changed; refusing to overwrite manual edits. "
            f"Expected {SOURCE_SHA256}, got {actual_hash}."
        )

    with zipfile.ZipFile(base_path, "r") as base_zip:
        base_slide = base_zip.read("ppt/slides/slide6.xml")
    with zipfile.ZipFile(emphasis_path, "r") as emphasis_zip:
        emphasis_slide = emphasis_zip.read("ppt/slides/slide6.xml")

    with zipfile.ZipFile(SOURCE, "r") as source_zip:
        slide_rels = source_zip.read("ppt/slides/_rels/slide6.xml.rels")
        replacements: dict[str, bytes] = {
            "ppt/slides/slide6.xml": base_slide,
        }

        presentation_xml = source_zip.read("ppt/presentation.xml")
        replacements["ppt/presentation.xml"] = replace_once(
            presentation_xml,
            b'<p:sldId id="261" r:id="rId13"/>',
            b'<p:sldId id="261" r:id="rId13"/><p:sldId id="288" r:id="rId40"/>',
            "Slide 7 insertion point",
        )

        presentation_rels = source_zip.read("ppt/_rels/presentation.xml.rels")
        relationship = (
            b'<Relationship Id="rId40" '
            b'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" '
            b'Target="slides/slide33.xml"/>'
        )
        replacements["ppt/_rels/presentation.xml.rels"] = replace_once(
            presentation_rels,
            b"</Relationships>",
            relationship + b"</Relationships>",
            "presentation relationship",
        )

        content_types = source_zip.read("[Content_Types].xml")
        override = (
            b'<Override PartName="/ppt/slides/slide33.xml" '
            b'ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'
        )
        replacements["[Content_Types].xml"] = replace_once(
            content_types,
            b"</Types>",
            override + b"</Types>",
            "content type",
        )

        additions = {
            "ppt/slides/slide33.xml": emphasis_slide,
            "ppt/slides/_rels/slide33.xml.rels": slide_rels,
        }

        temporary = OUTPUT.with_suffix(".tmp.pptx")
        if temporary.exists():
            temporary.unlink()
        with zipfile.ZipFile(temporary, "w") as output_zip:
            for info in source_zip.infolist():
                output_zip.writestr(
                    info, replacements.get(info.filename, source_zip.read(info.filename))
                )
            for name, payload in additions.items():
                output_zip.writestr(name, payload)
        temporary.replace(OUTPUT)


def main() -> None:
    with tempfile.TemporaryDirectory() as directory:
        directory_path = Path(directory)
        base_path = directory_path / "base.pptx"
        emphasis_path = directory_path / "emphasis.pptx"
        build_variants(base_path, emphasis_path)
        assemble(base_path, emphasis_path)
    print(f"Created {OUTPUT}")


if __name__ == "__main__":
    main()
