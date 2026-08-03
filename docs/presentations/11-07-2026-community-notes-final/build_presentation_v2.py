from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


HERE = Path(__file__).resolve().parent
DEFAULT_OUT = HERE / "community-notes-final-presentation-v2.pptx"
OUT = Path(os.environ.get("OUTPUT_PATH", DEFAULT_OUT))
PREVIEW_SLIDE = int(os.environ.get("PREVIEW_SLIDE", "0"))

W, H = 13.333, 7.5
FONT = "Helvetica Neue"

WHITE = "FFFFFF"
INK = "111318"
MUTED = "73777F"
LIGHT = "DADDE2"
FAINT = "ECEEF1"
BLUE = "3977F6"
CORAL = "FF705B"
GREEN = "1FA873"
GOLD = "D79A24"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
blank = prs.slide_layouts[6]


def new_slide():
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(WHITE)
    return slide


def textbox(slide, text, x, y, w, h, size=24, color=INK, bold=False,
            align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def bullet(slide, text, x, y, w, size=18, color=INK, dot=BLUE, bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.44))
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "• "; r.font.name = FONT; r.font.size = Pt(size)
    r.font.bold = True; r.font.color.rgb = rgb(dot)
    r = p.add_run(); r.text = text; r.font.name = FONT; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = rgb(color)
    return box


def line(slide, x1, y1, x2, y2, color=INK, width=1.5):
    shp = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    shp.line.color.rgb = rgb(color)
    shp.line.width = Pt(width)
    return shp


def circle(slide, x, y, d, fill=INK, outline=None, width=0.8):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid(); shp.fill.fore_color.rgb = rgb(fill)
    shp.line.color.rgb = rgb(outline or fill); shp.line.width = Pt(width)
    return shp


def outline_round_rect(slide, x, y, w, h, color=INK, width=1.2, radius=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.background()
    shp.line.color.rgb = rgb(color)
    shp.line.width = Pt(width)
    return shp


def title(slide, text, kicker, number):
    textbox(slide, kicker.upper(), 0.66, 0.46, 4.7, 0.25, 10, MUTED, True)
    textbox(slide, text, 0.66, 0.86, 11.4, 0.62, 30, INK, True)
    textbox(slide, f"{number:02d}", 12.12, 0.50, 0.48, 0.20, 9, MUTED, True, PP_ALIGN.RIGHT)


def source(slide, text):
    textbox(slide, f"• {text}", 0.68, 7.06, 11.9, 0.18, 8.2, MUTED)


def user(slide, x, y, scale=1.0, color=INK, active=False):
    head = 0.20 * scale
    circle(slide, x, y, head, color)
    line(slide, x + head/2, y + head, x + head/2, y + 0.53*scale, color, 2.4 if active else 1.5)
    line(slide, x - 0.04*scale, y + 0.34*scale, x + 0.24*scale, y + 0.34*scale,
         color, 2.4 if active else 1.5)
    line(slide, x + head/2, y + 0.53*scale, x - 0.03*scale, y + 0.74*scale,
         color, 2.4 if active else 1.5)
    line(slide, x + head/2, y + 0.53*scale, x + 0.23*scale, y + 0.74*scale,
         color, 2.4 if active else 1.5)


def post_note(slide, x, y, w=2.8, h=2.6, status=None, accent=GREEN):
    outline_round_rect(slide, x, y, w, h, LIGHT, 1.0)
    circle(slide, x + 0.20, y + 0.20, 0.24, INK)
    line(slide, x + 0.60, y + 0.28, x + 1.45, y + 0.28, INK, 1.5)
    line(slide, x + 0.20, y + 0.72, x + w - 0.25, y + 0.72, MUTED, 1.2)
    line(slide, x + 0.20, y + 0.98, x + w - 0.55, y + 0.98, MUTED, 1.2)
    line(slide, x + 0.20, y + 1.24, x + w - 0.35, y + 1.24, MUTED, 1.2)
    line(slide, x + 0.20, y + 1.55, x + w - 0.20, y + 1.55, LIGHT, 1.0)
    outline_round_rect(slide, x + 0.20, y + 1.78, w - 0.40, 0.55, LIGHT, 1.0)
    line(slide, x + 0.38, y + 1.96, x + w - 0.70, y + 1.96, accent, 1.4)
    line(slide, x + 0.38, y + 2.13, x + w - 1.03, y + 2.13, accent, 1.4)
    if status:
        circle(slide, x + w - 0.42, y + h - 0.38, 0.16, accent)
        textbox(slide, status, x + w - 1.28, y + h - 0.43, 0.74, 0.18,
                8.5, accent, True, PP_ALIGN.RIGHT)


def rating_marks(slide, x, y):
    circle(slide, x, y, 0.15, GREEN)
    line(slide, x + 0.03, y + 0.08, x + 0.07, y + 0.12, WHITE, 1.1)
    line(slide, x + 0.07, y + 0.12, x + 0.13, y + 0.04, WHITE, 1.1)
    circle(slide, x + 0.34, y, 0.15, WHITE, MUTED, 1.0)
    line(slide, x + 0.37, y + 0.075, x + 0.46, y + 0.075, MUTED, 1.0)


def crowd(slide, x, y, cols, rows, dx=0.38, dy=0.38, color=INK, d=0.10,
          active_indices=()):
    points = []
    idx = 0
    for r in range(rows):
        for c in range(cols):
            px = x + c*dx + (0.14 if r % 2 else 0)
            py = y + r*dy
            active = idx in active_indices
            circle(slide, px, py, d*(2.05 if active else 1.0),
                   CORAL if active else color)
            points.append((px + d/2, py + d/2, active))
            idx += 1
    return points


def slide_1():
    slide = new_slide()
    textbox(slide, "WHO SPEAKS\nFOR THE CROWD?", 0.72, 0.82, 6.2, 1.75, 38, INK, True)
    bullet(slide, "Cross-Constituency Aggregation", 0.77, 3.20, 5.6, 18, MUTED, BLUE)
    bullet(slide, "for Community Notes", 0.77, 3.68, 4.8, 18, MUTED, BLUE)
    points_a = crowd(slide, 8.25, 1.10, 5, 4, 0.38, 0.42, BLUE, 0.10)
    points_b = crowd(slide, 8.25, 4.72, 5, 3, 0.38, 0.42, CORAL, 0.10)
    outline_round_rect(slide, 11.15, 2.78, 1.35, 1.62, INK, 1.1)
    line(slide, 11.42, 3.25, 12.20, 3.25, INK, 1.4)
    line(slide, 11.42, 3.58, 12.03, 3.58, INK, 1.4)
    line(slide, points_a[-1][0], points_a[-1][1], 11.15, 3.22, BLUE, 1.5)
    line(slide, points_b[-1][0], points_b[-1][1], 11.15, 3.95, CORAL, 1.5)
    textbox(slide, "10 MIN", 0.77, 6.76, 0.8, 0.18, 9, MUTED, True)
    textbox(slide, "01", 12.12, 6.76, 0.48, 0.18, 9, MUTED, True, PP_ALIGN.RIGHT)


def slide_2():
    slide = new_slide(); title(slide, "One note. One decision.", "What is Community Notes?", 2)
    user(slide, 1.15, 2.62, 1.35, INK)
    post_note(slide, 4.05, 1.82, 3.05, 3.05)
    rating_marks(slide, 9.72, 3.06)
    line(slide, 2.05, 3.05, 3.55, 3.05, LIGHT, 1.4)
    textbox(slide, "→", 3.56, 2.84, 0.32, 0.30, 21, MUTED, True, PP_ALIGN.CENTER)
    line(slide, 7.52, 3.05, 9.32, 3.05, LIGHT, 1.4)
    textbox(slide, "→", 9.26, 2.84, 0.32, 0.30, 21, MUTED, True, PP_ALIGN.CENTER)
    bullet(slide, "Sees the note", 0.86, 5.55, 2.8, 17, INK, BLUE, True)
    bullet(slide, "Reads the context", 4.02, 5.55, 3.2, 17, INK, BLUE, True)
    bullet(slide, "Rates helpfulness", 9.35, 5.55, 3.1, 17, INK, GREEN, True)
    source(slide, "X Community Notes (2026)")


def slide_3():
    slide = new_slide(); title(slide, "Now multiply that by thousands", "The crowd", 3)
    user(slide, 0.95, 2.82, 1.05, INK)
    textbox(slide, "→", 2.15, 3.02, 0.40, 0.30, 21, MUTED, True, PP_ALIGN.CENTER)
    pts = crowd(slide, 3.10, 1.95, 15, 9, 0.42, 0.43, INK, 0.085)
    outline_round_rect(slide, 10.42, 2.37, 1.55, 2.02, INK, 1.1)
    line(slide, 10.72, 2.92, 11.65, 2.92, INK, 1.3)
    line(slide, 10.72, 3.27, 11.48, 3.27, INK, 1.3)
    line(slide, 10.72, 3.62, 11.58, 3.62, GREEN, 1.5)
    for idx in (11, 28, 47, 65, 82, 103, 122):
        px, py, _ = pts[idx]
        line(slide, px, py, 10.42, 3.38, LIGHT, 0.7)
    bullet(slide, "Same action", 0.90, 6.14, 2.5, 16, INK, BLUE)
    bullet(slide, "Thousands of raters", 4.20, 6.14, 3.4, 16, INK, BLUE)
    bullet(slide, "One visibility decision", 8.75, 6.14, 3.8, 16, INK, GREEN)
    source(slide, "X Community Notes (2026)")


def slide_4():
    slide = new_slide(); title(slide, "But the crowd is not equal", "The problem", 4)
    active = (5, 22, 58, 81, 104)
    pts = crowd(slide, 0.95, 1.90, 17, 10, 0.43, 0.43, INK, 0.075, active)
    note_x, note_y = 9.64, 2.48
    outline_round_rect(slide, note_x, note_y, 1.55, 2.00, INK, 1.1)
    line(slide, note_x + 0.28, note_y + 0.55, note_x + 1.25, note_y + 0.55, INK, 1.3)
    line(slide, note_x + 0.28, note_y + 0.93, note_x + 1.08, note_y + 0.93, INK, 1.3)
    line(slide, note_x + 0.28, note_y + 1.31, note_x + 1.17, note_y + 1.31, GREEN, 1.4)
    for idx, (px, py, act) in enumerate(pts):
        if act:
            for offset in (-0.06, 0.02, 0.10):
                line(slide, px, py + offset, note_x, note_y + 1.0, CORAL, 1.1)
        elif idx in (8, 41, 76, 137):
            line(slide, px, py, note_x, note_y + 1.0, LIGHT, 0.7)
    textbox(slide, "64", 11.56, 2.04, 1.05, 0.62, 35, CORAL, True, PP_ALIGN.RIGHT)
    textbox(slide, "• power raters", 11.12, 2.76, 1.48, 0.26, 12, MUTED, True, PP_ALIGN.RIGHT)
    textbox(slide, "11×", 11.40, 3.58, 1.20, 0.50, 29, INK, True, PP_ALIGN.RIGHT)
    textbox(slide, "• more active", 11.12, 4.18, 1.48, 0.26, 12, MUTED, True, PP_ALIGN.RIGHT)
    bullet(slide, "A few rate constantly", 1.00, 6.25, 3.2, 16, INK, CORAL)
    bullet(slide, "Most rate occasionally", 4.90, 6.25, 3.3, 16, INK, MUTED)
    source(slide, "Nudo et al. (2026); authors’ 200k analysis")


def slide_5():
    slide = new_slide(); title(slide, "The map inherits the imbalance", "Current bridging", 5)
    # Latent axis pulled upward by active users.
    line(slide, 0.95, 4.78, 11.85, 2.52, LIGHT, 2.0)
    for i in range(18):
        x = 1.10 + i*0.55
        y = 4.58 - i*0.11 + (0.34 if i % 3 == 0 else -0.08)
        circle(slide, x, y, 0.10, BLUE if i < 9 else CORAL)
    for x, y in ((7.85, 1.95), (8.35, 1.72), (8.72, 2.05)):
        circle(slide, x, y, 0.27, CORAL)
        line(slide, x + 0.13, y + 0.27, x + 0.05, 3.35, CORAL, 1.2)
    textbox(slide, "LATENT VIEWPOINT MAP", 0.98, 5.05, 3.1, 0.24, 10, MUTED, True)
    bullet(slide, "Activity shapes the map", 1.00, 6.05, 3.3, 17, INK, CORAL, True)
    bullet(slide, "The map defines ‘bridging’", 6.55, 6.05, 4.1, 17, INK, BLUE, True)
    source(slide, "X Community Notes (2026); Nudo et al. (2026)")


def slide_6():
    slide = new_slide(); title(slide, "Consult constituencies directly", "Our shift", 6)
    left = crowd(slide, 1.05, 2.00, 7, 7, 0.42, 0.45, BLUE, 0.10)
    right = crowd(slide, 8.65, 2.00, 7, 7, 0.42, 0.45, CORAL, 0.10)
    outline_round_rect(slide, 5.92, 2.52, 1.50, 1.92, INK, 1.1)
    line(slide, 6.20, 3.05, 7.15, 3.05, INK, 1.3)
    line(slide, 6.20, 3.40, 7.00, 3.40, INK, 1.3)
    line(slide, 6.20, 3.75, 7.10, 3.75, GREEN, 1.4)
    line(slide, left[-1][0], left[-1][1], 5.92, 3.18, BLUE, 1.5)
    line(slide, right[0][0], right[0][1], 7.42, 3.18, CORAL, 1.5)
    bullet(slide, "Recover groups", 1.05, 5.68, 2.5, 15, INK, BLUE)
    bullet(slide, "Ask each group", 5.14, 5.68, 2.6, 15, INK, GREEN)
    bullet(slide, "Aggregate explicitly", 9.26, 5.68, 3.0, 15, INK, CORAL)
    textbox(slide, "• Switzerland · Belgium · Bosnia · Northern Ireland",
            3.32, 6.35, 6.8, 0.24, 10.5, MUTED, False, PP_ALIGN.CENTER)
    source(slide, "Linder & Mueller (2021); Lijphart (1977); Bieber (2006); McGarry & O’Leary (2009)")


def slide_7():
    slide = new_slide(); title(slide, "Four design principles", "Philosophy", 7)
    xline = 1.42
    line(slide, xline, 1.87, xline, 6.27, LIGHT, 1.5)
    entries = [
        (2.00, "P1", "Presence", "Every group enters", BLUE),
        (3.06, "P2", "Non-compensation", "Rejection still matters", CORAL),
        (4.12, "P3", "Symmetry", "No default winner", GOLD),
        (5.18, "P4", "Behavioral recovery", "No outside labels", GREEN),
    ]
    for y, p, head, sub, col in entries:
        circle(slide, xline - 0.09, y + 0.08, 0.18, col)
        textbox(slide, p, 1.86, y - 0.02, 0.65, 0.26, 12, col, True)
        textbox(slide, head, 2.60, y - 0.10, 3.55, 0.34, 19, INK, True)
        bullet(slide, sub, 7.00, y - 0.06, 3.5, 15, MUTED, col)


def slide_8():
    slide = new_slide(); title(slide, "Principles become operations", "Implementation", 8)
    xline = 1.42
    line(slide, xline, 1.82, xline, 5.78, LIGHT, 1.5)
    entries = [
        (1.92, "P1", "≥3 ratings", "from every group", BLUE),
        (2.82, "P2", "Geometric mean", "a soft veto", CORAL),
        (3.72, "P3", "Same rule", "no size weighting", GOLD),
        (4.62, "P4", "200k raters", "Method-B recovery", GREEN),
    ]
    for y, p, head, sub, col in entries:
        circle(slide, xline - 0.09, y + 0.06, 0.18, col)
        textbox(slide, p, 1.86, y - 0.03, 0.60, 0.24, 11, col, True)
        textbox(slide, head, 2.58, y - 0.08, 2.95, 0.30, 17, INK, True)
        bullet(slide, sub, 5.60, y - 0.05, 2.55, 13.5, MUTED, col)
    textbox(slide, "90%", 8.58, 2.15, 1.00, 0.42, 24, BLUE, True)
    textbox(slide, "×", 9.58, 2.17, 0.34, 0.34, 20, MUTED, True, PP_ALIGN.CENTER)
    textbox(slide, "10%", 9.98, 2.15, 1.00, 0.42, 24, CORAL, True)
    textbox(slide, "→", 11.00, 2.17, 0.36, 0.34, 20, MUTED, True, PP_ALIGN.CENTER)
    textbox(slide, "30%", 11.45, 2.10, 1.00, 0.46, 28, GREEN, True, PP_ALIGN.RIGHT)
    bullet(slide, "not 50%", 9.52, 3.08, 2.4, 14, MUTED, GREEN)
    line(slide, 8.62, 4.15, 12.30, 4.15, LIGHT, 1.2)
    textbox(slide, "100k notes  →  200k raters  →  > .5  →  Gabriel",
            8.58, 4.47, 3.85, 0.36, 13, INK, True, PP_ALIGN.RIGHT)
    source(slide, "Authors’ 200k Representative pipeline")


def slide_9():
    slide = new_slide(); title(slide, "From hidden to validated", "Results", 9)
    y = 3.35
    xs = [0.92, 4.15, 7.45, 10.50]
    vals = [("44,722", "representative picks", INK),
            ("6,832", "shown · 15%", BLUE),
            ("13,655", "CCA candidates", CORAL),
            ("3,896", "validated", GREEN)]
    line(slide, 1.35, y, 11.70, y, LIGHT, 1.8)
    for x, (big, small, col) in zip(xs, vals):
        circle(slide, x + 0.58, y - 0.09, 0.18, col)
        textbox(slide, big, x, 2.16, 1.45, 0.52, 28, col, True, PP_ALIGN.CENTER)
        bullet(slide, small, x - 0.05, 3.78, 1.95, 12.5, MUTED, col)
    textbox(slide, "→", 3.25, 3.06, 0.36, 0.32, 18, MUTED, True, PP_ALIGN.CENTER)
    textbox(slide, "+", 6.58, 3.05, 0.30, 0.32, 18, MUTED, True, PP_ALIGN.CENTER)
    textbox(slide, "→", 9.70, 3.06, 0.36, 0.32, 18, MUTED, True, PP_ALIGN.CENTER)
    bullet(slide, "46% potential before validation", 1.00, 5.30, 4.4, 17, INK, BLUE, True)
    bullet(slide, "3,896 pass Gabriel", 7.45, 5.30, 3.5, 17, INK, GREEN, True)
    source(slide, "Authors’ results · two camps · behavioral ≠ demographic · model-based validation")


def slide_10():
    slide = new_slide(); title(slide, "An aggregation problem", "Takeaway", 10)
    left = crowd(slide, 0.95, 2.00, 7, 7, 0.42, 0.45, BLUE, 0.095)
    right = crowd(slide, 8.70, 2.00, 7, 7, 0.42, 0.45, CORAL, 0.095)
    outline_round_rect(slide, 5.98, 2.50, 1.42, 1.90, GREEN, 1.4)
    line(slide, 6.25, 3.05, 7.12, 3.05, INK, 1.2)
    line(slide, 6.25, 3.40, 6.98, 3.40, INK, 1.2)
    line(slide, 6.25, 3.75, 7.08, 3.75, GREEN, 1.5)
    line(slide, left[-1][0], left[-1][1], 5.98, 3.18, BLUE, 1.5)
    line(slide, right[0][0], right[0][1], 7.40, 3.18, CORAL, 1.5)
    bullet(slide, "Make groups explicit", 0.98, 5.72, 3.0, 15.5, INK, BLUE)
    bullet(slide, "Require cross-group support", 4.70, 5.72, 3.8, 15.5, INK, GREEN)
    bullet(slide, "Validate content separately", 9.08, 5.72, 3.5, 15.5, INK, CORAL)
    textbox(slide, "PRESERVE DISAGREEMENT BEFORE AGGREGATING IT",
            3.03, 6.52, 7.35, 0.24, 10, MUTED, True, PP_ALIGN.CENTER)


SLIDES = [slide_1, slide_2, slide_3, slide_4, slide_5,
          slide_6, slide_7, slide_8, slide_9, slide_10]

if PREVIEW_SLIDE:
    if PREVIEW_SLIDE < 1 or PREVIEW_SLIDE > len(SLIDES):
        raise SystemExit("PREVIEW_SLIDE must be between 1 and 10")
    SLIDES[PREVIEW_SLIDE - 1]()
else:
    for build in SLIDES:
        build()

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"Saved {OUT}")
print(f"Slides: {len(prs.slides)}")
