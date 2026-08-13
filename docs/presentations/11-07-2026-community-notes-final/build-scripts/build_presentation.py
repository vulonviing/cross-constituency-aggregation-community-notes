from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


HERE = Path(__file__).resolve().parent
OUT = HERE / "community-notes-final-presentation.pptx"

W = 13.333
H = 7.5

BG = "F7F7F5"
INK = "111318"
MUTED = "6E737B"
LINE = "E2E3E0"
WHITE = "FFFFFF"
BLUE = "2F6BFF"
BLUE_SOFT = "E8EFFF"
CORAL = "FF654F"
CORAL_SOFT = "FFEAE5"
GREEN = "22A06B"
GREEN_SOFT = "E4F5ED"
GOLD = "F2B84B"
GOLD_SOFT = "FFF3D6"
GRAY_TILE = "ECEDE9"

FONT = "Helvetica Neue"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
blank = prs.slide_layouts[6]


def shape(slide, kind, x, y, w, h, fill, line=None, radius=True):
    shp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.line.color.rgb = rgb(line or fill)
    shp.line.width = Pt(0.8)
    return shp


def rounded(slide, x, y, w, h, fill=WHITE, line=None):
    return shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line)


def circle(slide, x, y, d, fill, line=None):
    return shape(slide, MSO_SHAPE.OVAL, x, y, d, d, fill, line)


def line(slide, x1, y1, x2, y2, color=INK, width=2, dash=None):
    shp = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    shp.line.color.rgb = rgb(color)
    shp.line.width = Pt(width)
    if dash is not None:
        shp.line.dash_style = dash
    return shp


def textbox(slide, text, x, y, w, h, size=24, color=INK, bold=False,
            align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0,
            font=FONT, tracking=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def bullets(slide, items, x, y, w, h, size=22, color=INK, gap=10,
            bullet_color=None, bold_first=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.0
        p.alignment = PP_ALIGN.LEFT
        r1 = p.add_run()
        r1.text = "• "
        r1.font.name = FONT
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = rgb(bullet_color or color)
        r2 = p.add_run()
        r2.text = item
        r2.font.name = FONT
        r2.font.size = Pt(size)
        r2.font.bold = bold_first
        r2.font.color.rgb = rgb(color)
    return box


def title(slide, text, kicker=None, number=None):
    if kicker:
        textbox(slide, kicker.upper(), 0.78, 0.45, 5.5, 0.28, 10, MUTED, True)
    textbox(slide, text, 0.78, 0.78, 11.6, 0.72, 30, INK, True)
    if number is not None:
        textbox(slide, f"{number:02d}", 12.15, 0.5, 0.4, 0.25, 9, MUTED, True,
                align=PP_ALIGN.RIGHT)


def source(slide, text):
    textbox(slide, f"• {text}", 0.78, 7.08, 11.8, 0.2, 8.5, MUTED)


def add_note_icon(slide, x, y, s, fill=WHITE, stroke=INK):
    rounded(slide, x, y, s, s, fill, fill)
    rounded(slide, x + 0.21*s, y + 0.18*s, 0.58*s, 0.62*s, WHITE, LINE)
    line(slide, x + 0.31*s, y + 0.34*s, x + 0.69*s, y + 0.34*s, stroke, 1.4)
    line(slide, x + 0.31*s, y + 0.46*s, x + 0.61*s, y + 0.46*s, stroke, 1.4)
    line(slide, x + 0.31*s, y + 0.58*s, x + 0.65*s, y + 0.58*s, stroke, 1.4)
    circle(slide, x + 0.60*s, y + 0.66*s, 0.10*s, GREEN)


def add_people_icon(slide, x, y, s, a=BLUE, b=CORAL, bg=WHITE):
    rounded(slide, x, y, s, s, bg, bg)
    points = [
        (0.24, 0.30, a, 0.13), (0.43, 0.22, a, 0.15),
        (0.62, 0.31, b, 0.13), (0.72, 0.50, b, 0.15),
        (0.49, 0.48, a, 0.17), (0.30, 0.60, a, 0.12),
        (0.58, 0.68, b, 0.13),
    ]
    for px, py, col, d in points:
        circle(slide, x + px*s, y + py*s, d*s, col)


def add_bridge_icon(slide, x, y, s):
    rounded(slide, x, y, s, s, WHITE, WHITE)
    for i in range(3):
        circle(slide, x + (0.17 + i*0.10)*s, y + (0.66 - i*0.04)*s, 0.07*s, BLUE)
        circle(slide, x + (0.68 + i*0.08)*s, y + (0.58 + i*0.04)*s, 0.07*s, CORAL)
    line(slide, x + 0.33*s, y + 0.62*s, x + 0.70*s, y + 0.55*s, INK, 4)
    line(slide, x + 0.38*s, y + 0.62*s, x + 0.47*s, y + 0.37*s, INK, 2)
    line(slide, x + 0.47*s, y + 0.37*s, x + 0.64*s, y + 0.56*s, INK, 2)


def add_eye_icon(slide, x, y, s):
    rounded(slide, x, y, s, s, GOLD_SOFT, GOLD_SOFT)
    eye = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.18*s), Inches(y + 0.30*s),
                                 Inches(0.64*s), Inches(0.38*s))
    eye.fill.solid(); eye.fill.fore_color.rgb = rgb(WHITE)
    eye.line.color.rgb = rgb(INK); eye.line.width = Pt(1.5)
    circle(slide, x + 0.43*s, y + 0.42*s, 0.14*s, INK)


def add_scale_icon(slide, x, y, s):
    rounded(slide, x, y, s, s, CORAL_SOFT, CORAL_SOFT)
    line(slide, x + 0.50*s, y + 0.22*s, x + 0.50*s, y + 0.74*s, INK, 2)
    line(slide, x + 0.25*s, y + 0.38*s, x + 0.75*s, y + 0.49*s, INK, 3)
    line(slide, x + 0.25*s, y + 0.38*s, x + 0.17*s, y + 0.62*s, MUTED, 1.2)
    line(slide, x + 0.75*s, y + 0.49*s, x + 0.83*s, y + 0.65*s, MUTED, 1.2)
    circle(slide, x + 0.13*s, y + 0.60*s, 0.13*s, BLUE)
    for off in (0.71, 0.79):
        circle(slide, x + off*s, y + 0.63*s, 0.10*s, CORAL)


def country_tile(slide, x, y, code, label, fill, ink=INK):
    rounded(slide, x, y, 2.55, 1.75, fill, fill)
    rounded(slide, x + 0.18, y + 0.18, 0.62, 0.62, WHITE, WHITE)
    textbox(slide, code, x + 0.18, y + 0.34, 0.62, 0.22, 15, ink, True,
            PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    textbox(slide, label, x + 0.22, y + 1.03, 2.05, 0.38, 15, INK, True)


def principle_card(slide, x, y, num, label, icon, fill, accent):
    rounded(slide, x, y, 2.78, 2.05, fill, fill)
    rounded(slide, x + 0.19, y + 0.18, 0.58, 0.58, WHITE, WHITE)
    textbox(slide, num, x + 0.19, y + 0.35, 0.58, 0.18, 12, accent, True,
            PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    textbox(slide, label, x + 0.22, y + 1.02, 2.30, 0.36, 17, INK, True)
    textbox(slide, f"• {icon}", x + 0.22, y + 1.45, 2.25, 0.34, 11.5, MUTED)


# Slide 1
slide = prs.slides.add_slide(blank)
slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(BG)
textbox(slide, "WHO SPEAKS\nFOR THE CROWD?", 0.78, 0.72, 6.1, 2.1, 38, INK, True)
textbox(slide, "• Cross-Constituency Aggregation\n• for Community Notes",
        0.82, 3.05, 4.8, 0.9, 19, MUTED)
add_people_icon(slide, 7.40, 0.78, 2.35, BLUE, BLUE, BLUE_SOFT)
add_note_icon(slide, 9.93, 2.12, 2.35, WHITE, INK)
add_people_icon(slide, 7.40, 4.46, 2.35, CORAL, CORAL, CORAL_SOFT)
line(slide, 8.65, 2.95, 9.94, 3.18, BLUE, 2.5)
line(slide, 8.65, 5.06, 9.94, 3.92, CORAL, 2.5)
textbox(slide, "10 MIN", 0.82, 6.72, 1.0, 0.2, 9, MUTED, True)
textbox(slide, "01", 12.20, 6.72, 0.35, 0.2, 9, MUTED, True, PP_ALIGN.RIGHT)


# Slide 2
slide = prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(BG)
title(slide, "What is Community Notes?", "The system", 2)
tiles = [(0.95, "01", "Write context", BLUE_SOFT, BLUE),
         (4.25, "02", "Rate helpfulness", GRAY_TILE, INK),
         (7.55, "03", "Decide visibility", GREEN_SOFT, GREEN)]
for x, n, lab, fill, accent in tiles:
    rounded(slide, x, 2.05, 2.75, 2.75, fill, fill)
    rounded(slide, x + 0.22, 2.28, 0.66, 0.66, WHITE, WHITE)
    textbox(slide, n, x + 0.22, 2.49, 0.66, 0.18, 12, accent, True,
            PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    textbox(slide, f"• {lab}", x + 0.27, 4.10, 2.22, 0.36, 16, INK, True)
line(slide, 3.75, 3.43, 4.15, 3.43, MUTED, 1.7)
line(slide, 7.05, 3.43, 7.45, 3.43, MUTED, 1.7)
rounded(slide, 10.78, 2.05, 1.62, 2.75, INK, INK)
textbox(slide, "VISIBLE", 10.97, 3.16, 1.24, 0.28, 14, WHITE, True,
        PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
bullets(slide, ["The crowd decides which context appears"], 0.98, 5.62, 8.4, 0.5, 21, INK, 0, BLUE)
source(slide, "X Community Notes (2026)")


# Slide 3
slide = prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(BG)
title(slide, "Not just majority rule", "Bridging", 3)
rounded(slide, 0.95, 1.86, 5.35, 4.55, GRAY_TILE, GRAY_TILE)
rounded(slide, 7.03, 1.86, 5.35, 4.55, WHITE, WHITE)
add_people_icon(slide, 1.55, 2.34, 2.05, BLUE, BLUE, BLUE_SOFT)
add_note_icon(slide, 4.15, 2.34, 1.55, WHITE)
line(slide, 3.38, 3.37, 4.12, 3.15, BLUE, 3)
textbox(slide, "HIDDEN", 2.01, 5.45, 3.3, 0.3, 16, MUTED, True, PP_ALIGN.CENTER)
add_people_icon(slide, 7.55, 2.34, 1.80, BLUE, BLUE, BLUE_SOFT)
add_note_icon(slide, 9.74, 2.55, 1.35, GREEN_SOFT)
add_people_icon(slide, 11.36, 2.34, 1.80, CORAL, CORAL, CORAL_SOFT)
line(slide, 9.25, 3.25, 9.74, 3.20, BLUE, 3)
line(slide, 11.36, 3.25, 11.08, 3.20, CORAL, 3)
textbox(slide, "VISIBLE", 8.85, 5.45, 3.3, 0.3, 16, GREEN, True, PP_ALIGN.CENTER)
bullets(slide, ["Reward support that crosses viewpoints", "Estimate the bridge through latent patterns"],
        0.98, 6.58, 11.4, 0.45, 13, MUTED, 2, BLUE)
source(slide, "Wojcik et al. (2022); X Community Notes (2026)")


# Slide 4
slide = prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(BG)
title(slide, "The bridge is indirect", "The problem", 4)
add_eye_icon(slide, 0.98, 1.90, 3.00)
add_scale_icon(slide, 4.25, 1.90, 3.00)
rounded(slide, 7.52, 1.90, 4.82, 3.00, INK, INK)
textbox(slide, "64", 7.95, 2.30, 1.35, 0.85, 46, WHITE, True)
textbox(slide, "POWER\nRATERS", 9.58, 2.38, 2.10, 0.72, 14, WHITE, True)
textbox(slide, "11×", 7.95, 3.56, 1.35, 0.55, 25, CORAL, True)
textbox(slide, "MORE ACTIVE", 9.58, 3.68, 2.10, 0.28, 12, WHITE, True)
bullets(slide, ["Infer who voted", "Let activity shape the map", "Mistake the map for the crowd"],
        1.04, 5.45, 10.8, 0.98, 18, INK, 7, CORAL)
source(slide, "X Community Notes (2026); Nudo et al. (2026); authors’ 200k analysis")


# Slide 5
slide = prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(BG)
title(slide, "Consult constituencies directly", "Our shift", 5)
bullets(slide, ["Recover groups", "Measure approval", "Aggregate explicitly"],
        0.98, 1.76, 4.2, 1.4, 18, INK, 8, BLUE)
country_tile(slide, 0.98, 4.05, "CH", "Double majority", BLUE_SOFT, BLUE)
country_tile(slide, 3.78, 4.05, "BE", "Parallel consent", GOLD_SOFT, GOLD)
country_tile(slide, 6.58, 4.05, "BA", "Community veto", CORAL_SOFT, CORAL)
country_tile(slide, 9.38, 4.05, "NI", "Cross-support", GREEN_SOFT, GREEN)
rounded(slide, 6.28, 1.68, 6.12, 1.72, WHITE, WHITE)
textbox(slide, "“", 6.55, 1.77, 0.46, 0.42, 38, BLUE, True)
textbox(slide, "• Which rule should decide visibility?", 7.04, 2.20, 4.72, 0.42, 19, INK, True)
source(slide, "Linder & Mueller (2021); Lijphart (1977); Bieber (2006); McGarry & O’Leary (2009)")


# Slide 6
slide = prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(BG)
title(slide, "Four design principles", "Philosophy", 6)
principle_card(slide, 0.95, 1.86, "P1", "Presence", "Every group enters", BLUE_SOFT, BLUE)
principle_card(slide, 3.94, 1.86, "P2", "Non-compensation", "Rejection still matters", CORAL_SOFT, CORAL)
principle_card(slide, 6.93, 1.86, "P3", "Symmetry", "No default winner", GOLD_SOFT, GOLD)
principle_card(slide, 9.92, 1.86, "P4", "Behavioral recovery", "No outside labels", GREEN_SOFT, GREEN)
bullets(slide, ["Preserve disagreement before aggregating it"], 0.98, 5.55, 10.8, 0.5, 20, INK, 0, BLUE)


# Slide 7
slide = prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(BG)
title(slide, "Principles become operations", "Implementation", 7)
rows = [
    ("P1", "≥3 ratings", "from every group", BLUE_SOFT, BLUE),
    ("P2", "Geometric mean", "a soft veto", CORAL_SOFT, CORAL),
    ("P3", "Same rule", "no size weighting", GOLD_SOFT, GOLD),
    ("P4", "200k raters", "Method-B recovery", GREEN_SOFT, GREEN),
]
for i, (num, head, sub, fill, accent) in enumerate(rows):
    x = 0.95 + i * 3.00
    rounded(slide, x, 1.85, 2.75, 2.30, fill, fill)
    textbox(slide, num, x + 0.22, 2.10, 0.55, 0.26, 12, accent, True)
    textbox(slide, head, x + 0.22, 2.70, 2.20, 0.38, 18, INK, True)
    textbox(slide, f"• {sub}", x + 0.22, 3.36, 2.20, 0.30, 12, MUTED)
rounded(slide, 2.08, 4.83, 9.18, 1.18, INK, INK)
textbox(slide, "90%", 2.42, 5.20, 1.02, 0.30, 21, BLUE, True)
textbox(slide, "×", 3.52, 5.20, 0.30, 0.30, 19, WHITE, True, PP_ALIGN.CENTER)
textbox(slide, "10%", 3.94, 5.20, 1.02, 0.30, 21, CORAL, True)
textbox(slide, "→", 5.12, 5.20, 0.42, 0.30, 19, WHITE, True, PP_ALIGN.CENTER)
textbox(slide, "30%", 5.76, 5.14, 1.20, 0.38, 26, GREEN, True)
textbox(slide, "• not 50%", 7.22, 5.22, 1.55, 0.28, 14, WHITE, True)
textbox(slide, "C = √(s₁ × s₂)", 9.06, 5.22, 1.80, 0.28, 14, WHITE, True, PP_ALIGN.RIGHT)
source(slide, "Authors’ 200k Representative pipeline")


# Slide 8
slide = prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(BG)
title(slide, "What we tested", "Pipeline", 8)
pipeline = [
    ("100k", "notes", BLUE_SOFT, BLUE),
    ("200k", "raters", CORAL_SOFT, CORAL),
    ("2", "groups", GOLD_SOFT, GOLD),
    ("> .5", "candidate", GREEN_SOFT, GREEN),
    ("✓", "Gabriel", INK, WHITE),
]
for i, (big, small, fill, accent) in enumerate(pipeline):
    x = 0.78 + i * 2.53
    rounded(slide, x, 2.12, 2.12, 2.12, fill, fill)
    textbox(slide, big, x + 0.18, 2.73, 1.76, 0.48, 28, accent, True, PP_ALIGN.CENTER)
    textbox(slide, f"• {small}", x + 0.20, 3.52, 1.72, 0.28, 13, accent, True, PP_ALIGN.CENTER)
    if i < len(pipeline) - 1:
        line(slide, x + 2.17, 3.18, x + 2.43, 3.18, MUTED, 1.5)
bullets(slide, ["Recover", "Score", "Select", "Validate"], 1.04, 5.28, 10.9, 0.44, 17, INK, 0, BLUE)
source(slide, "Authors’ analysis of Community Notes data")


# Slide 9
slide = prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(BG)
title(slide, "Candidates, then validation", "Results", 9)
rounded(slide, 0.95, 1.86, 5.40, 4.45, WHITE, WHITE)
textbox(slide, "15%", 1.35, 2.31, 1.45, 0.74, 42, BLUE, True)
textbox(slide, "→", 2.96, 2.48, 0.55, 0.40, 25, MUTED, True, PP_ALIGN.CENTER)
textbox(slide, "46%", 3.72, 2.31, 1.58, 0.74, 42, GREEN, True)
textbox(slide, "• pre-validation visibility", 1.40, 3.40, 3.95, 0.30, 15, MUTED)
rounded(slide, 1.35, 4.35, 4.50, 0.50, GRAY_TILE, GRAY_TILE)
rounded(slide, 1.35, 4.35, 0.68, 0.50, BLUE, BLUE)
rounded(slide, 2.03, 4.35, 1.40, 0.50, GREEN, GREEN)
textbox(slide, "• +31 percentage points", 1.40, 5.24, 3.95, 0.32, 15, INK, True)

rounded(slide, 6.63, 1.86, 5.75, 4.45, INK, INK)
textbox(slide, "13,655", 7.08, 2.28, 2.20, 0.55, 31, WHITE, True)
textbox(slide, "• candidates", 9.44, 2.45, 1.70, 0.28, 14, WHITE, True)
line(slide, 8.20, 3.10, 8.20, 3.75, MUTED, 2)
textbox(slide, "3,896", 7.08, 4.04, 2.20, 0.58, 34, GREEN, True)
textbox(slide, "• validated", 9.44, 4.24, 1.70, 0.28, 14, WHITE, True)
textbox(slide, "• sourced\n• factual\n• rescue-worthy", 7.10, 5.02, 3.60, 0.80, 12.5, WHITE)
source(slide, "Authors’ pipeline results; Gabriel validation · two camps · model-based validation")


# Slide 10
slide = prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(INK)
textbox(slide, "TAKEAWAY", 0.82, 0.62, 2.0, 0.25, 10, "A7ABB2", True)
textbox(slide, "NOT ONLY A\nPREDICTION PROBLEM.", 0.82, 1.22, 6.4, 1.52, 31, WHITE, True)
textbox(slide, "AN AGGREGATION\nPROBLEM.", 0.82, 2.94, 6.4, 1.52, 31, GREEN, True)
bullets(slide, ["Make constituencies explicit", "Require support across them", "Validate content separately"],
        0.88, 5.28, 5.6, 1.12, 17, WHITE, 7, GREEN)
add_people_icon(slide, 7.35, 0.92, 2.15, BLUE, BLUE, "182239")
add_people_icon(slide, 9.94, 4.44, 2.15, CORAL, CORAL, "321C1A")
add_note_icon(slide, 9.03, 2.47, 2.15, GREEN_SOFT, INK)
line(slide, 8.58, 2.53, 9.24, 3.08, BLUE, 2.5)
line(slide, 10.50, 4.44, 10.36, 4.15, CORAL, 2.5)
textbox(slide, "10", 12.20, 6.72, 0.35, 0.2, 9, "A7ABB2", True, PP_ALIGN.RIGHT)


prs.save(OUT)
print(f"Saved {OUT}")
print(f"Slides: {len(prs.slides)}")
