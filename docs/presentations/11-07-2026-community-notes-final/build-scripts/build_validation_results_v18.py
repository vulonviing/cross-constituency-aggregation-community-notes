from __future__ import annotations

import hashlib
import tempfile
import zipfile
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from build_two_problems_v4 import (
    BLUE,
    CORAL,
    GREEN,
    INK,
    LIGHT,
    MUTED,
    H,
    W,
    base_slide,
    bullet,
    circle,
    line,
    rgb,
    source,
    textbox,
    title,
)


HERE = Path(__file__).resolve().parent
PROJECT_ROOT = HERE.parents[2]
SOURCE = HERE / "community-notes-final-presentation-v17.pptx"
OUTPUT = HERE / "community-notes-final-presentation-v18.pptx"
SOURCE_SHA256 = "52dba63868f9acdbff045145feb6973e7579ce8ec2fcfd3b65ba1f2fa5db0e01"
TARGET_SLIDE = "ppt/slides/slide33.xml"

STAGE1_DIR = (
    PROJECT_ROOT
    / "data"
    / "llm_validation"
    / "runs"
    / "gemma-4-31b-it-scckn-v1"
)
STAGE15_DIR = (
    PROJECT_ROOT
    / "data"
    / "llm_validation"
    / "runs"
    / "gemma-4-31b-it-scckn-stage1-5-opinion-v1"
)
STAGE2_DIR = (
    PROJECT_ROOT
    / "data"
    / "llm_validation"
    / "runs"
    / "gemma-4-31b-it-scckn-stage2-expanded-v1"
)


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def validate_counts() -> None:
    stage1 = pd.read_parquet(STAGE1_DIR / "stage1_results.parquet")
    stage15 = pd.read_parquet(STAGE15_DIR / "stage1_5_results.parquet")
    stage2 = pd.read_parquet(STAGE2_DIR / "stage2_results.parquet")

    expected_stage1 = {
        "sourced_factual_information": 10_096,
        "opinion_or_speculation": 1_703,
        "irrelevant_trivial_or_spam": 1_340,
        "unsourced_context_or_claim": 373,
        "hostile_troll_or_derogatory": 142,
    }
    actual_stage1 = stage1["final_label"].value_counts().to_dict()
    unresolved = int(stage1["final_label"].isna().sum())
    if len(stage1) != 13_655 or actual_stage1 != expected_stage1 or unresolved != 1:
        raise RuntimeError(
            f"Stage 1 counts changed: rows={len(stage1)}, "
            f"labels={actual_stage1}, unresolved={unresolved}"
        )

    expected_stage15 = {
        "sourced_factual_core_absent": 1_423,
        "sourced_factual_core_present": 280,
    }
    actual_stage15 = stage15["stage1_5_label"].value_counts().to_dict()
    if len(stage15) != 1_703 or actual_stage15 != expected_stage15:
        raise RuntimeError(
            f"Stage 1.5 counts changed: rows={len(stage15)}, labels={actual_stage15}"
        )

    route_pass = (
        stage2.groupby(["admission_route", "passes_rescue_threshold"])
        .size()
        .to_dict()
    )
    expected_route_pass = {
        ("stage1_5_recall", False): 249,
        ("stage1_5_recall", True): 31,
        ("strict_stage1", False): 1_569,
        ("strict_stage1", True): 8_527,
    }
    if len(stage2) != 10_376 or route_pass != expected_route_pass:
        raise RuntimeError(
            f"Stage 2 counts changed: rows={len(stage2)}, routes={route_pass}"
        )

    content_stop = 1_423 + 1_340 + 373 + 142 + 1
    stage2_stop = int((~stage2["passes_rescue_threshold"]).sum())
    final_rescues = int(stage2["passes_rescue_threshold"].sum())
    if (content_stop, stage2_stop, final_rescues) != (3_279, 1_818, 8_558):
        raise RuntimeError(
            "Derived validation totals changed: "
            f"content_stop={content_stop}, stage2_stop={stage2_stop}, "
            f"final_rescues={final_rescues}"
        )


def compact_node(
    slide,
    x: float,
    number: str,
    label: str,
    detail: str,
    color: str,
) -> None:
    textbox(slide, number, x, 1.73, 1.65, 0.50, 28, color, True, PP_ALIGN.CENTER)
    circle(slide, x + 0.72, 2.46, 0.20, color)
    bullet(
        slide,
        label,
        x - 0.18,
        2.84,
        2.05,
        11.6,
        MUTED,
        color,
        True,
        PP_ALIGN.CENTER,
    )
    textbox(
        slide,
        detail,
        x - 0.28,
        3.18,
        2.22,
        0.20,
        8.8,
        color,
        True,
        PP_ALIGN.CENTER,
    )


def reason_row(slide, x: float, y: float, width: float, count: str, label: str) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.24))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = False
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.space_after = Pt(0)

    number_run = paragraph.add_run()
    number_run.text = f"{count}  "
    number_run.font.name = "Helvetica Neue"
    number_run.font.size = Pt(10.2)
    number_run.font.bold = True
    number_run.font.color.rgb = rgb(CORAL)

    label_run = paragraph.add_run()
    label_run.text = label
    label_run.font.name = "Helvetica Neue"
    label_run.font.size = Pt(10.2)
    label_run.font.color.rgb = rgb(MUTED)


def validation_results(slide) -> None:
    title(slide, "RESULTS · VALIDATION", "From candidates to validated rescues", "13")

    # Main canonical Gemma funnel.
    line(slide, 1.56, 2.56, 9.02, 2.56, LIGHT, 1.8)
    compact_node(slide, 0.82, "13,655", "CCA candidates", "frozen universe", INK)
    textbox(slide, "→", 3.58, 2.34, 0.36, 0.30, 18, MUTED, True,
            PP_ALIGN.CENTER)
    compact_node(
        slide,
        4.08,
        "10,376",
        "Stage 2 judgments",
        "10,096 strict + 280 recall",
        BLUE,
    )
    textbox(slide, "→", 6.82, 2.34, 0.36, 0.30, 18, MUTED, True,
            PP_ALIGN.CENTER)
    compact_node(
        slide,
        7.34,
        "8,558",
        "Final rescues · ≥50",
        "8,527 strict + 31 recall",
        GREEN,
    )

    # Content exclusions after the targeted Stage 1.5 recall pass.
    line(slide, 2.66, 2.66, 2.66, 3.45, CORAL, 1.0)
    textbox(slide, "3,279  CONTENT STOP", 1.18, 3.52, 3.10, 0.25,
            13.0, CORAL, True, PP_ALIGN.CENTER)
    content_reasons = [
        ("1,423", "opinion / factual core absent"),
        ("1,340", "irrelevant / trivial / spam"),
        ("373", "unsourced context / claim"),
        ("142", "hostile / derogatory"),
        ("1", "unresolved"),
    ]
    for index, (count, label) in enumerate(content_reasons):
        reason_row(slide, 1.48, 3.91 + index * 0.25, 3.15, count, label)

    # Stage 2 scores below the preregistered rescue threshold.
    line(slide, 6.04, 2.66, 6.04, 3.45, CORAL, 1.0)
    textbox(slide, "1,818  BELOW 50", 5.00, 3.52, 2.08, 0.25,
            13.0, CORAL, True, PP_ALIGN.CENTER)
    reason_row(slide, 5.13, 4.10, 2.55, "1,818", "rescue-worthiness < 50")

    # Completed production-run strip replaces the historical rerun messaging.
    line(slide, 0.82, 5.42, 12.50, 5.42, LIGHT, 1.0)
    textbox(slide, "GEMMA 4 31B IT", 0.94, 5.70, 2.35, 0.24,
            11.5, GREEN, True)
    textbox(slide, "COMPLETED CANONICAL RUN", 3.45, 5.70, 2.60, 0.24,
            11.0, INK, True)
    bullet(slide, "8,527 strict route", 6.55, 5.68, 2.25,
           11.3, MUTED, BLUE, True)
    bullet(slide, "31 recall route", 9.25, 5.68, 2.10,
           11.3, MUTED, GREEN, True)
    textbox(
        slide,
        "10,376 / 10,376 complete · zero unresolved Stage 2 judgments",
        3.20,
        6.28,
        6.95,
        0.24,
        11.5,
        GREEN,
        True,
        PP_ALIGN.CENTER,
    )
    source(slide, "Canonical Gemma 4 31B IT run · SCCKN · BF16 · threshold ≥50")


def build_generated(path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(W)
    presentation.slide_height = Inches(H)
    validation_results(base_slide(presentation))
    presentation.save(path)


def assemble(generated_path: Path) -> None:
    actual_hash = sha256(SOURCE)
    if actual_hash != SOURCE_SHA256:
        raise SystemExit(
            "V17 changed; refusing to overwrite manual edits. "
            f"Expected {SOURCE_SHA256}, got {actual_hash}."
        )

    with zipfile.ZipFile(generated_path, "r") as generated_zip:
        replacement_slide = generated_zip.read("ppt/slides/slide1.xml")

    with zipfile.ZipFile(SOURCE, "r") as source_zip:
        temporary = OUTPUT.with_suffix(".tmp.pptx")
        if temporary.exists():
            temporary.unlink()
        with zipfile.ZipFile(temporary, "w") as output_zip:
            for info in source_zip.infolist():
                payload = (
                    replacement_slide
                    if info.filename == TARGET_SLIDE
                    else source_zip.read(info.filename)
                )
                output_zip.writestr(info, payload)
        temporary.replace(OUTPUT)


def main() -> None:
    validate_counts()
    with tempfile.TemporaryDirectory() as directory:
        generated_path = Path(directory) / "validation-results.pptx"
        build_generated(generated_path)
        assemble(generated_path)
    print(f"Created {OUTPUT}")


if __name__ == "__main__":
    main()
