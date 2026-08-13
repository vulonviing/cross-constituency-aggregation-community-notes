from __future__ import annotations

import hashlib
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from build_two_problems_v4 import (
    BLUE,
    CORAL,
    GREEN,
    INK,
    MUTED,
    H,
    W,
    base_slide,
    rgb,
)
from build_validation_results_v18 import validate_counts
from build_validation_share_v19 import validation_results_with_share


HERE = Path(__file__).resolve().parent
SOURCE = HERE / "community-notes-final-presentation-v21.pptx"
OUTPUT = HERE / "community-notes-final-presentation-v22.pptx"
SOURCE_SHA256 = "e1cb21cfd4c63ea6f790bd92c2521e022141166200d2f340c0576d777e421c06"
TARGET_SLIDE = "ppt/slides/slide33.xml"


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def score_row(
    slide,
    y: float,
    score: str,
    label: str,
    detail: str,
    color: str,
) -> None:
    box = slide.shapes.add_textbox(Inches(5.05), Inches(y), Inches(3.95), Inches(0.21))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = False
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.space_after = Pt(0)

    dot = paragraph.add_run()
    dot.text = "•  "
    dot.font.name = "Helvetica Neue"
    dot.font.size = Pt(8.6)
    dot.font.bold = True
    dot.font.color.rgb = rgb(color)

    score_run = paragraph.add_run()
    score_run.text = f"{score}  "
    score_run.font.name = "Helvetica Neue"
    score_run.font.size = Pt(8.6)
    score_run.font.bold = True
    score_run.font.color.rgb = rgb(color)

    label_run = paragraph.add_run()
    label_run.text = f"{label} · "
    label_run.font.name = "Helvetica Neue"
    label_run.font.size = Pt(8.6)
    label_run.font.bold = True
    label_run.font.color.rgb = rgb(INK if color != CORAL else CORAL)

    detail_run = paragraph.add_run()
    detail_run.text = detail
    detail_run.font.name = "Helvetica Neue"
    detail_run.font.size = Pt(8.6)
    detail_run.font.color.rgb = rgb(MUTED)


def validation_results_with_bands(slide) -> None:
    validation_results_with_share(slide)

    old_rows = [
        shape
        for shape in slide.shapes
        if hasattr(shape, "text") and "rescue-worthiness < 50" in shape.text
    ]
    if len(old_rows) != 1:
        raise RuntimeError(f"Expected one old Stage 2 stop row, found {len(old_rows)}")
    old_row = old_rows[0]._element
    old_row.getparent().remove(old_row)

    rows = [
        ("90–100", "outstanding", "clear, traceable, self-contained", GREEN),
        ("70–89", "strong", "well-supported and useful", GREEN),
        ("40–69", "mixed", "useful elements, material gaps", BLUE),
        ("10–39", "weak", "limited support or clarity", CORAL),
        ("0–9", "minimal", "little usable justification", CORAL),
    ]
    for index, row in enumerate(rows):
        score_row(slide, 3.88 + index * 0.22, *row)
    score_row(slide, 5.02, "≥50", "validated", "display-quality threshold", GREEN)


def build_generated(path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(W)
    presentation.slide_height = Inches(H)
    validation_results_with_bands(base_slide(presentation))
    presentation.save(path)


def assemble(generated_path: Path) -> None:
    actual_hash = sha256(SOURCE)
    if actual_hash != SOURCE_SHA256:
        raise SystemExit(
            "V21 changed; refusing to overwrite manual edits. "
            f"Expected {SOURCE_SHA256}, got {actual_hash}."
        )

    with zipfile.ZipFile(generated_path, "r") as generated_zip:
        replacement_slide = generated_zip.read("ppt/slides/slide1.xml")

    with zipfile.ZipFile(SOURCE, "r") as source_zip:
        temporary = OUTPUT.with_suffix(".tmp.pptx")
        if temporary.exists():
            temporary.unlink()
        with zipfile.ZipFile(temporary, "w") as output_zip:
            for info in source_zip.infolist():
                payload = (
                    replacement_slide
                    if info.filename == TARGET_SLIDE
                    else source_zip.read(info.filename)
                )
                output_zip.writestr(info, payload)
        temporary.replace(OUTPUT)


def main() -> None:
    validate_counts()
    with tempfile.TemporaryDirectory() as directory:
        generated_path = Path(directory) / "validation-bands.pptx"
        build_generated(generated_path)
        assemble(generated_path)
    print(f"Created {OUTPUT}")


if __name__ == "__main__":
    main()
