from __future__ import annotations

import hashlib
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from build_results_v9 import wip
from build_two_problems_v4 import (
    BLUE,
    CORAL,
    GREEN,
    INK,
    LIGHT,
    MUTED,
    W,
    H,
    base_slide,
    bullet,
    circle,
    line,
    rgb,
    source,
    textbox,
    title,
)


HERE = Path(__file__).resolve().parent
SOURCE = HERE / "community-notes-final-presentation-v11.pptx"
OUTPUT = HERE / "community-notes-final-presentation-v12.pptx"
SOURCE_SHA256 = "0f0684eb53a3d7fb4afeafa8f986da1fe10f4e37d99f98718287e5bf91165f61"


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def compact_node(slide, x: float, number: str, label: str, color: str) -> None:
    textbox(slide, number, x, 1.73, 1.65, 0.50, 28, color, True, PP_ALIGN.CENTER)
    circle(slide, x + 0.72, 2.46, 0.20, color)
    bullet(slide, label, x - 0.18, 2.86, 2.05, 11.8, MUTED, color, True,
           PP_ALIGN.CENTER)


def reason_row(slide, x: float, y: float, width: float, count: str, label: str) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.24))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = False
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.space_after = Pt(0)

    number_run = paragraph.add_run()
    number_run.text = f"{count}  "
    number_run.font.name = "Helvetica Neue"
    number_run.font.size = Pt(10.5)
    number_run.font.bold = True
    number_run.font.color.rgb = rgb(CORAL)

    label_run = paragraph.add_run()
    label_run.text = label
    label_run.font.name = "Helvetica Neue"
    label_run.font.size = Pt(10.5)
    label_run.font.color.rgb = rgb(MUTED)


def validation_reasons(slide) -> None:
    title(slide, "RESULTS · VALIDATION", "From candidates to validated rescues", "13")
    wip(slide)

    # Compact top funnel leaves enough room for explicit stop reasons.
    line(slide, 1.56, 2.56, 9.02, 2.56, LIGHT, 1.8)
    compact_node(slide, 0.82, "13,655", "CCA candidates", INK)
    textbox(slide, "→", 3.58, 2.34, 0.36, 0.30, 18, MUTED, True,
            PP_ALIGN.CENTER)
    compact_node(slide, 4.08, "8,051", "Stage 1 pass", BLUE)
    textbox(slide, "→", 6.82, 2.34, 0.36, 0.30, 18, MUTED, True,
            PP_ALIGN.CENTER)
    compact_node(slide, 7.34, "3,896", "Stage 2 pass · ≥50", GREEN)

    # Stage 1: five mutually exclusive historical stop categories.
    line(slide, 5.00, 2.66, 5.00, 3.31, CORAL, 1.1)
    textbox(slide, "5,604  STAGE 1 STOP", 3.30, 3.32, 3.45, 0.25,
            13.2, CORAL, True, PP_ALIGN.CENTER)
    stage1_reasons = [
        ("3,610", "opinion / speculation"),
        ("687", "irrelevant / trivial / spam"),
        ("616", "hostile / troll / derogatory"),
        ("474", "unsourced context / claim"),
        ("217", "unresolved / strict ties"),
    ]
    for index, (count, label) in enumerate(stage1_reasons):
        reason_row(slide, 3.43, 3.73 + index * 0.28, 3.38, count, label)

    # Stage 2: sourced notes that do not clear the rescue-worthiness threshold.
    line(slide, 8.24, 2.66, 8.24, 3.31, CORAL, 1.1)
    textbox(slide, "4,155  STAGE 2 STOP", 7.03, 3.32, 2.48, 0.25,
            13.2, CORAL, True, PP_ALIGN.CENTER)
    reason_row(slide, 7.12, 3.92, 2.82, "4,155", "rescue-worthiness < 50")

    # Preserve the existing model-rerun strip in full.
    line(slide, 0.82, 5.42, 12.50, 5.42, LIGHT, 1.0)
    textbox(slide, "GPT-4o-mini baseline", 0.94, 5.70, 2.30, 0.24,
            11.5, MUTED, True)
    textbox(slide, "→", 3.23, 5.68, 0.36, 0.26,
            16, MUTED, True, PP_ALIGN.CENTER)
    textbox(slide, "MiMo v2.5 Pro rerun", 3.70, 5.70, 2.45, 0.24,
            11.5, BLUE, True)
    bullet(slide, "higher reasoning capacity", 6.55, 5.68, 2.55,
           11.5, MUTED, BLUE)
    bullet(slide, "accuracy expected to improve", 9.30, 5.68, 2.82,
           11.5, MUTED, GREEN)
    textbox(slide, "FINAL COUNTS WILL CHANGE", 4.50, 6.30, 4.35, 0.25,
            12, CORAL, True, PP_ALIGN.CENTER)
    source(slide, "Historical Gabriel outputs · MiMo rerun underway")


def build_generated(path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(W)
    presentation.slide_height = Inches(H)
    validation_reasons(base_slide(presentation))
    presentation.save(path)


def assemble(generated_path: Path) -> None:
    actual_hash = sha256(SOURCE)
    if actual_hash != SOURCE_SHA256:
        raise SystemExit(
            "V11 changed; refusing to overwrite manual edits. "
            f"Expected {SOURCE_SHA256}, got {actual_hash}."
        )

    with zipfile.ZipFile(generated_path, "r") as generated_zip:
        replacement_slide = generated_zip.read("ppt/slides/slide1.xml")

    with zipfile.ZipFile(SOURCE, "r") as source_zip:
        temporary = OUTPUT.with_suffix(".tmp.pptx")
        if temporary.exists():
            temporary.unlink()
        with zipfile.ZipFile(temporary, "w") as output_zip:
            for info in source_zip.infolist():
                payload = (
                    replacement_slide
                    if info.filename == "ppt/slides/slide31.xml"
                    else source_zip.read(info.filename)
                )
                output_zip.writestr(info, payload)
        temporary.replace(OUTPUT)


def main() -> None:
    with tempfile.TemporaryDirectory() as directory:
        generated_path = Path(directory) / "validation-reasons.pptx"
        build_generated(generated_path)
        assemble(generated_path)
    print(f"Created {OUTPUT}")


if __name__ == "__main__":
    main()
