from __future__ import annotations

import hashlib
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


HERE = Path(__file__).resolve().parent
SOURCE = HERE / "community-notes-final-presentation-v2.pptx"
OUTPUT = HERE / "community-notes-final-presentation-v3.pptx"
SOURCE_SHA256 = "7f1aa1bb21f8aad811e03076cd624c5fee79a5a949617fed98149c5b1030b992"
TARGET_SLIDE = "ppt/slides/slide8.xml"

W, H = 13.333, 7.5
FONT = "Helvetica Neue"
WHITE = "FFFFFF"
INK = "111318"
MUTED = "73777F"
LIGHT = "DADDE2"
BLUE = "3977F6"
CORAL = "FF705B"
GREEN = "1FA873"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def textbox(slide, text, x, y, w, h, size=24, color=INK, bold=False,
            align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def bullet(slide, text, x, y, w, size=14, color=INK, dot=GREEN, bold=False,
           align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.38))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    r1 = paragraph.add_run()
    r1.text = "• "
    r1.font.name = FONT
    r1.font.size = Pt(size)
    r1.font.bold = True
    r1.font.color.rgb = rgb(dot)
    r2 = paragraph.add_run()
    r2.text = text
    r2.font.name = FONT
    r2.font.size = Pt(size)
    r2.font.bold = bold
    r2.font.color.rgb = rgb(color)
    return box


def line(slide, x1, y1, x2, y2, color=INK, width=1.5):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    connector.line.color.rgb = rgb(color)
    connector.line.width = Pt(width)
    return connector


def circle(slide, x, y, diameter, fill=INK, outline=None, width=0.8):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x), Inches(y), Inches(diameter), Inches(diameter),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(outline or fill)
    shape.line.width = Pt(width)
    return shape


def title(slide):
    textbox(slide, "CURRENT BRIDGING", 0.66, 0.46, 4.8, 0.25, 10, MUTED, True)
    textbox(slide, "Activity moves the bridge", 0.66, 0.86, 11.5, 0.62, 30, INK, True)
    textbox(slide, "05", 12.12, 0.50, 0.48, 0.20, 9, MUTED, True, PP_ALIGN.RIGHT)


def user_row(slide, x, y, color, count=5, active=False):
    for index in range(count):
        px = x + index * 0.34
        circle(slide, px, y, 0.10, color)
        if active:
            # Repeated rating ticks: same users, more activity—not more people.
            for tick in range(4):
                ty = y - 0.16 - tick * 0.12
                line(slide, px + 0.05, ty, px + 0.05, ty + 0.055, color, 1.2)


def comparison_axis(slide, x1, x2, y, marker_x, marker_label):
    line(slide, x1, y, x2, y, LIGHT, 1.8)
    circle(slide, x1 - 0.04, y - 0.04, 0.08, BLUE)
    circle(slide, x2 - 0.04, y - 0.04, 0.08, CORAL)
    line(slide, marker_x, y - 0.33, marker_x, y + 0.31, GREEN, 2.0)
    circle(slide, marker_x - 0.07, y - 0.07, 0.14, GREEN)
    textbox(slide, marker_label, marker_x - 0.62, y + 0.45, 1.24, 0.23,
            9, GREEN, True, PP_ALIGN.CENTER)


def build_replacement_pptx(path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(W)
    presentation.slide_height = Inches(H)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(WHITE)

    title(slide)

    # Two scenes, separated only by a hairline—no cards or panels.
    line(slide, 6.67, 1.72, 6.67, 4.92, LIGHT, 1.0)
    textbox(slide, "BEFORE · EQUAL ACTIVITY", 0.82, 1.65, 3.5, 0.25, 10, MUTED, True)
    textbox(slide, "AFTER · CORAL RATES 5× MORE", 7.05, 1.65, 4.0, 0.25, 10, MUTED, True)

    # Before: identical rating volume.
    user_row(slide, 1.08, 2.26, BLUE)
    user_row(slide, 4.35, 2.26, CORAL)
    textbox(slide, "100 × 10% = 10", 0.86, 2.67, 2.40, 0.30, 14, BLUE, True)
    textbox(slide, "100 × 90% = 90", 4.00, 2.67, 2.25, 0.30, 14, CORAL, True, PP_ALIGN.RIGHT)
    comparison_axis(slide, 0.98, 6.15, 3.45, 3.565, "50% SIGNAL")
    textbox(slide, "(10 + 90) / 200 = 50%", 1.63, 4.28, 3.92, 0.34,
            19, INK, True, PP_ALIGN.CENTER)

    # After: same users and within-group approval; coral emits 5x rating volume.
    user_row(slide, 7.32, 2.26, BLUE)
    user_row(slide, 10.58, 2.26, CORAL, active=True)
    textbox(slide, "100 × 10% = 10", 7.10, 2.67, 2.40, 0.30, 14, BLUE, True)
    textbox(slide, "500 × 90% = 450", 10.18, 2.67, 2.40, 0.30, 14, CORAL, True, PP_ALIGN.RIGHT)
    after_x1, after_x2 = 7.22, 12.40
    after_marker = after_x1 + (after_x2 - after_x1) * 0.77
    comparison_axis(slide, after_x1, after_x2, 3.45, after_marker, "77% SIGNAL")
    textbox(slide, "(10 + 450) / 600 ≈ 77%", 7.68, 4.28, 4.25, 0.34,
            19, INK, True, PP_ALIGN.CENTER)

    # Full-width CCA punchline.
    line(slide, 0.82, 5.10, 12.52, 5.10, LIGHT, 1.1)
    bullet(slide, "Same within-group support: 10% | 90%", 0.88, 5.46, 3.85,
           13, MUTED, BLUE, True)
    textbox(slide, "√(.10 × .90) = 30%", 4.58, 5.35, 4.18, 0.52,
            27, GREEN, True, PP_ALIGN.CENTER)
    bullet(slide, "unchanged by rating volume", 9.12, 5.46, 3.08,
           13, MUTED, GREEN, True, PP_ALIGN.RIGHT)
    textbox(slide, "CCA STAYS SYMMETRIC", 5.35, 6.02, 2.65, 0.22,
            10, GREEN, True, PP_ALIGN.CENTER)

    textbox(slide, "• Toy illustration—not the Community Notes scoring formula.",
            0.68, 6.75, 5.8, 0.18, 8.2, MUTED)
    textbox(slide, "• X Community Notes (2026); Nudo et al. (2026)",
            7.08, 6.75, 5.45, 0.18, 8.2, MUTED, False, PP_ALIGN.RIGHT)

    presentation.save(path)


def replace_slide_xml(replacement_pptx: Path) -> None:
    actual_hash = sha256(SOURCE)
    if actual_hash != SOURCE_SHA256:
        raise SystemExit(
            "Source V2 changed; refusing to overwrite manual edits. "
            f"Expected {SOURCE_SHA256}, got {actual_hash}."
        )

    with zipfile.ZipFile(replacement_pptx, "r") as replacement_zip:
        replacement_xml = replacement_zip.read("ppt/slides/slide1.xml")

    temp_output = OUTPUT.with_suffix(".tmp.pptx")
    if temp_output.exists():
        temp_output.unlink()

    with zipfile.ZipFile(SOURCE, "r") as source_zip, zipfile.ZipFile(temp_output, "w") as output_zip:
        for info in source_zip.infolist():
            data = replacement_xml if info.filename == TARGET_SLIDE else source_zip.read(info.filename)
            output_zip.writestr(info, data)

    temp_output.replace(OUTPUT)


def main() -> None:
    with tempfile.TemporaryDirectory() as tmpdir:
        replacement_pptx = Path(tmpdir) / "replacement-slide.pptx"
        build_replacement_pptx(replacement_pptx)
        replace_slide_xml(replacement_pptx)
    print(f"Source V2 preserved: {sha256(SOURCE)}")
    print(f"Saved V3: {OUTPUT}")


if __name__ == "__main__":
    main()
